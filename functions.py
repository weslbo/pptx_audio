import requests
import markdownify
import re
import os
import azure.cognitiveservices.speech as speechsdk
from pptx import Presentation
from lxml import etree
from bs4 import BeautifulSoup
from readability import Document
from openai import AzureOpenAI

def removeNotes(input_pptx_file, output_pptx_file):
    presentation = Presentation(input_pptx_file)
    for slide in presentation.slides:
        if slide.notes_slide is not None:
            slide.notes_slide.notes_text_frame.clear()
    presentation.save(output_pptx_file)
    print("Notes removed from " + output_pptx_file)

def retrieveLearnContent(input_pptx_file, output_pptx_file):
    presentation = Presentation(input_pptx_file)

    iSlide = 1
    for slide in presentation.slides:
        title = "" if slide.shapes.title is None else slide.shapes.title.text
        print("Slide : " + str(iSlide).zfill(2) + " (" + title + ")")

        if slide.notes_slide is not None:
            notes = slide.notes_slide.notes_text_frame.text

            if notes.startswith("http"):
                url = notes.split('\n')[0]

                print("- Retrieving page url: " + url)
                
                response = requests.get(url)
                doc = Document(response.content)

                markdown = markdownify.markdownify(str(doc.summary()), heading_style="ATX")
                markdown = re.sub('\n{3,}', '\n\n', markdown)
                markdown = markdown.replace("[Continue](/en-us/)", "")
                
                slide.notes_slide.notes_text_frame.text = markdown
        else:
            print("- No URL found in slide notes")

        iSlide += 1
            
    presentation.save(output_pptx_file)
    print("Finished retrieving learn content on " + output_pptx_file)

def generateSSML(input_pptx_file, output_pptx_file, slideid_togenerate=-1):
    openai_client = AzureOpenAI(azure_endpoint=os.getenv("AZURE_OPENAI_API_ENDPOINT"), api_version="2023-07-01-preview", api_key=os.getenv("AZURE_OPENAI_API_KEY"))
    presentation = Presentation(input_pptx_file)

    iSlide = 1
    for slide in presentation.slides:
        title = "" if slide.shapes.title is None else slide.shapes.title.text
        print("Slide : " + str(iSlide).zfill(2) + " (" + title + ")")

        if slideid_togenerate == iSlide or slideid_togenerate == -1:
            if slide.notes_slide is not None and slide.notes_slide.notes_text_frame.text != "":    
                print("- Generating conversation for slide")

                notes = slide.notes_slide.notes_text_frame.text

                # message_text = [
                #     {"role":"system","content":"""
                #         You are a training instructor [Andrew] and you will be teaching this topic, based on the content below.
                #         - Create a transcript that is engaging and informative.
                #         - Students will listen to this talk, so the transcript should use natural spoken language.
                #         - The transcript should be about the topic in the notes, do not talk about any other topic.
                #         - You can use a bit of humor, but make sure it is appropriate for a professional setting.
                        
                #         Transform the text to the Speech Syntheses Markup Language (SSML).  
                #         - [Andrew] should use the voice with name en-US-AndrewNeural
                #         - There is no need for introductions anymore. No "welcome" needed.
                #         - Immediately jump into the topic.
                #         - The output should be XML.
                #         - Make sure that every line is wrapped between <voice> and </voice> element.                
                #         - Finally, make sure there is the following element at the start: <speak xmlns""http://www.w3.org/2001/10/synthesis"" xmlns:mstts=""http://www.w3.org/2001/mstts"" xmlns:emo=""http://www.w3.org/2009/10/emotionml"" version=""1.0"" xml:lang=""en-US"">
                #         - End the XML document with the following element: </speak>
                #         - Delete [Andrew]: from the transcript.
                #         ------------""" + notes},
                #     {"role":"user","content":"generate the transcript"}
                # ]

                message_text = [
                    {"role":"system","content":"""
                        You're going to create a transcript for an engaging conversation between [Brian] and [Andrew], based on the content below. They should ask each other questions and respond to each other.
                        Do not talk about any other topic.
                        Transform the text to the Speech Syntheses Markup Language (SSML).  
                        - [Brian] should use the voice with name en-US-BrianNeural
                        - [Andrew] should use the voice with name en-US-AndrewNeural
                        - There is no need for introductions anymore. No "welcome" needed.
                        - The output should be XML.
                        - Make sure that every line is wrapped between <voice> and </voice> element.                
                        - Finally, make sure there is the following element at the start: <speak xmlns""http://www.w3.org/2001/10/synthesis"" xmlns:mstts=""http://www.w3.org/2001/mstts"" xmlns:emo=""http://www.w3.org/2009/10/emotionml"" version=""1.0"" xml:lang=""en-US"">
                        - End the XML document with the following element: </speak>
                        - Delete [Brian]: and [Andrew]: from the transcript.
                        ------------
                     """ + notes},
                    {"role":"user","content":"generate the conversation"}
                ]

                completion = openai_client.chat.completions.create(
                    model="gpt-35-turbo-16k",
                    messages = message_text,
                    temperature=0.8,
                    max_tokens=8000,
                    top_p=0.95,
                    frequency_penalty=0,
                    presence_penalty=0,
                    stop=None
                )

                output = completion.choices[0].message.content

                slide.notes_slide.notes_text_frame.text = output
            else:
                print("- No notes found on slide")
        else:
            print("- Skipping slide")

        iSlide += 1
                                        
    presentation.save(output_pptx_file)
    print("Finished generated SSML for " + output_pptx_file)

def generateAudio(input_pptx_file, output_pptx_file, prefix, slideid_togenerate=-1):
    presentation = Presentation(input_pptx_file)

    iSlide = 1
    for slide in presentation.slides:
        title = "" if slide.shapes.title is None else slide.shapes.title.text
        print("Slide : " + str(iSlide).zfill(2) + " (" + title + ")")

        if slideid_togenerate == iSlide or slideid_togenerate == -1:
            if slide.notes_slide is not None and slide.notes_slide.notes_text_frame.text != "":
                print("- Generating audio for slide")

                notes = slide.notes_slide.notes_text_frame.text

                audio_file = f"{prefix}-{iSlide}.mp3"
                generate_audio_fragment(notes, audio_file)
                audio = slide.shapes.add_movie(audio_file, 0, 0, 1, 1)

                # Adjust auto play of audio
                # tree = audio._element.getparent().getparent().getnext().getnext().getnext()
                # timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond'][0]
                # print(timing)
                # timing.set('delay', '0')
        else:
            print("- Skipping slide")
        
        iSlide += 1
            

    presentation.save(output_pptx_file)
    print("Finished generated Audio for " + output_pptx_file)


def generate_audio_fragment(ssml, savelocation):
    print(f"- Creating audio {savelocation}")
    
    service_region = "eastus"
    speech_key = os.getenv("SPEECH_API_KEY")
    speech_config = speechsdk.SpeechConfig(subscription=speech_key, region=service_region)
    speech_config.set_speech_synthesis_output_format(speechsdk.SpeechSynthesisOutputFormat.Audio24Khz96KBitRateMonoMp3)  

    file_config = speechsdk.audio.AudioOutputConfig(filename=savelocation)
    speech_synthesizer = speechsdk.SpeechSynthesizer(speech_config=speech_config, audio_config=file_config)  

    result = speech_synthesizer.speak_ssml_async(ssml).get()
    return result

def listAllSlides(pptx_file):
    presentation = Presentation(pptx_file)
    all_slides = []
    for slide in presentation.slides:
        all_slides.append({ 'id': slide.slide_id, 'title': "" if slide.shapes.title is None else slide.shapes.title.text, 'url': "" })

    return all_slides

def addAudioFragment(pptx_file, slide_number, audio_file):
    presentation = Presentation(pptx_file)
    slide = presentation.slides[slide_number]
    audio = slide.shapes.add_movie(audio_file, 0, 0, 1, 1)

    # Adjust auto play of audio
    tree = audio._element.getparent().getparent().getnext().getnext().getnext()
    timing = [el for el in tree.iterdescendants() if etree.QName(el).localname == 'cond'][0]
    print(timing)
    timing.set('delay', '0')

    presentation.save(pptx_file)
    print("Audio added to slide " + str(slide_number) + " in " + pptx_file)