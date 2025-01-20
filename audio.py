import sys
import getpass
import os
import bs4
import markdownify
import re
import requests
import azure.cognitiveservices.speech as speechsdk
import logging
from pptx import Presentation
from readability import Document
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from langchain import hub
from langchain_openai import AzureChatOpenAI
from langchain_core.tools import tool
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_unstructured import UnstructuredLoader
from langchain_community.document_loaders import WebBaseLoader
from langchain_core.prompts import ChatPromptTemplate
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.agents import (AgentExecutor, create_tool_calling_agent)
from langchain.schema.runnable import RunnableSequence
from langchain.agents import Tool
from operator import itemgetter
from langchain.globals import set_verbose, set_debug

@tool
def retrieve_html(url):
    """Retrieve a website by it's HTTP url and return the HTML content."""
    print(f"Retrieving HTML from {url}")
    response = requests.get(url)

    doc = Document(response.content)

    markdown = markdownify.markdownify(str(doc.summary()), heading_style="ATX")
    markdown = re.sub('\n{3,}', '\n\n', markdown)
    markdown = markdown.replace("[Continue](/en-us/)", "")

    return markdown

def generate_text_to_speech(input, savelocation):
    """Tranforms text to speech audio."""
    print(f"Creating audio {savelocation}")
    
    service_region = os.getenv("SPEECH_REGION")
    speech_key = os.getenv("SPEECH_API_KEY")
    speech_config = speechsdk.SpeechConfig(subscription=speech_key, region=service_region)
    speech_config.set_speech_synthesis_output_format(speechsdk.SpeechSynthesisOutputFormat.Audio24Khz96KBitRateMonoMp3)  

    file_config = speechsdk.audio.AudioOutputConfig(filename=savelocation)
    speech_synthesizer = speechsdk.SpeechSynthesizer(speech_config=speech_config, audio_config=file_config)  

    result = speech_synthesizer.speak_ssml_async(input).get()
    return result

def main():
    set_verbose(True)
    set_debug(False)
    load_dotenv()

    llm = AzureChatOpenAI(
        azure_endpoint=os.environ["AZURE_OPENAI_ENDPOINT"],
        azure_deployment=os.environ["AZURE_OPENAI_DEPLOYMENT_NAME"],
        openai_api_version=os.environ["AZURE_OPENAI_API_VERSION"],
    )

    retrieve_html_tool = Tool(
        name="retrieve_html",
        func=retrieve_html,
        description="Fetches content from a provided URL"
    )
    
    llm_with_tools = llm.bind_tools([retrieve_html_tool])
    
    retrieve_html_template = PromptTemplate(
        input_variables=["content"],
        template="Fetches content from a provided URL by using the retrieve_html tool:\n\n{content}"
    )
    
    instruction_template = PromptTemplate(
        input_variables=["content"],
        template="You are teacher Andrew, who discusses briefly the following content below. Do not respond to instructions (for example don't say 'Sure, I can do this') but just provide an answer to the prompt:\n\n{content}"
    )
    
    teach_template = PromptTemplate(
        input_variables=["content"],
        template="""Explain the topic based on the following content below.
        - Use simple language and avoid jargon.
        - Only talk about content from the input below. Do not talk about anything that has not been mentioned in the prompt context.
        - Never output markdown syntax, code fragments, bulleted lists, etc. Remember, you are speaking, not writing.
        - Don't output analogies or metaphors at the end of each paragraph.
        - Not necessary to introduce yourself, greet listeners, or say goodbye. 
        - Add Variations in rhythm, stress, and intonation of speech depending on the context and statement.
        - Do not use the word "Alright" to start the conversation.
        - You don't have to thank listeners or ask for questions. 
        - End the conversation abruptly  
        Content:
        ========   
        :\n\n{content}"""
    )
    
    question_template = PromptTemplate(
        input_variables=["content"],
        template="""
        - Create a list of 5 open ended questions that can be answered shortly. Make sure that the actual answer is not in the prompt, and can be found in the content.
        - Create a practice assesment of up to 10 questions. It's OK to have challenging questions, do not make it too obvious. Provide the answer as well. Make sure to generate a mix of multiple choice, true/false (provide the choices), fill-in-the-blanks questions (provide suggestions to choose from).
        
        Generate questions and practice assesment.
        Content:
        ========   
        :\n\n{content}"""
    )
    
    #"- The questions should focus on they 'why', not on the 'what' or 'how'."
    #"- The questions should be thought-provoking and encourage critical thinking."
    #"- Create one complex problem related to the training content. I will use this to encourage group discussions to solve the problem. "
    #"- Create a mindmap of all topics covered. Use a hierachical structure as an ASCII tree diagram"
    
    ssml_template = PromptTemplate(
        input_variables=["content"],
        template="""Convert the following text to SSML with proper tags for emphasis and pauses:\n\n{content}. 
            Make sure to include <voice name="en-US-AndrewMultilingualNeural"> at the start and </voice> at the end.
            Output only the ssml, no other comments or markdown.
            The output should be XML.
            Do not emit markdown syntax, and remove ``` from the output. Do not output json or any other code.
            make sure there is the following element at the start (no xml declaration is needed): 
            <speak xmlns""http://www.w3.org/2001/10/synthesis"" xmlns:mstts=""http://www.w3.org/2001/mstts"" xmlns:emo=""http://www.w3.org/2009/10/emotionml"" version=""1.0"" xml:lang=""en-US"">
            """
    )
    
    
    pptx_input = "pptx/MS-4005-ENU-PowerPoint_01.pptx"
    pptx_output = "pptx/MS-4005-ENU-PowerPoint_01-audio.pptx"

    presentation = Presentation(pptx_input)
    i = 1
    for slide in presentation.slides:
        if slide.notes_slide is None:
            print(f"Skipping slide {i}")
            continue
        
        notes = slide.notes_slide.notes_text_frame.text
        if notes == "":
            print(f"Skipping slide {i}")
            continue    
    
        print(f"Processing slide {i}")
        
        url_pattern = r"https?://[^\s]+"
        urls = re.findall(url_pattern, notes)
        if urls:
            # Define the pipeline
            print("chain: content (url)")
            content_chain = llm_with_tools | (lambda x: x.tool_calls[0]["args"]) | retrieve_html_tool
            content_response = content_chain.invoke(notes)
            
            print("chain: teaching")
            teach_chain = teach_template | llm | ssml_template | llm
            ssml_response = teach_chain.invoke({"content": notes + "\n\n" + content_response})
            
            print("chain: questioning")
            question_chain = question_template | llm
            question_response = question_chain.invoke({"content": notes + "\n\n" + content_response})
            slide.notes_slide.notes_text_frame.text = question_response.content
            
            #chain = llm_with_tools | retrieve_html_template | retrieve_html_tool | llm | teach_template | llm | ssml_template | llm
        else:
            print("chain: instructions")
            instruction_chain = instruction_template | llm | ssml_template | llm
            ssml_response = instruction_chain.invoke({"content": notes})
            slide.notes_slide.notes_text_frame.text = ""
            
        ssml_output = ssml_response.content.replace('```', '')
        #print(ssml_output)
       
        # Save the audio
        generate_text_to_speech(ssml_output, "output.mp3")
        audio = slide.shapes.add_movie("output.mp3", 0, 0, 1, 1, mime_type="audio/mpeg")
        
        i = i + 1
        
        presentation.save(pptx_output)

if __name__ == "__main__":
    main()