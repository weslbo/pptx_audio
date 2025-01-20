import sys
import getpass
import os
import bs4
import markdownify
import re
import requests
import azure.cognitiveservices.speech as speechsdk
import logging
import uuid
import json
import time
from pptx import Presentation
from pptx.util import Inches, Cm
from readability import Document
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from langchain import hub
from langchain_openai import AzureChatOpenAI
from langchain_core.tools import tool
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_unstructured import UnstructuredLoader
from langchain_community.document_loaders import WebBaseLoader
from langchain_core.prompts import ChatPromptTemplate
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.agents import (AgentExecutor, create_tool_calling_agent)
from langchain.schema.runnable import RunnableSequence
from langchain.agents import Tool
from operator import itemgetter
from langchain.globals import set_verbose, set_debug
from openai import AzureOpenAI
from dotenv import load_dotenv
from urllib.parse import urlparse, parse_qs

load_dotenv()

SPEECH_ENDPOINT = f"https://{os.getenv('SPEECH_REGION')}.api.cognitive.microsoft.com/" 
API_VERSION = "2024-04-15-preview"
SUBSCRIPTION_KEY = os.getenv("SPEECH_API_KEY")

@tool
def retrieve_html(url):
    """Retrieve a website by it's HTTP url and return the HTML content."""
    print(f"Retrieving HTML from {url}")
    response = requests.get(url)

    doc = Document(response.content)

    markdown = markdownify.markdownify(str(doc.summary()), heading_style="ATX")
    markdown = re.sub('\n{3,}', '\n\n', markdown)
    markdown = markdown.replace("[Continue](/en-us/)", "")

    return markdown


def generate_video(transcript: str):
    job_id = str(uuid.uuid4())
    download_url = None
    if submit_synthesis(job_id, transcript):
        while True:
            status = get_synthesis(job_id)
            if status == 'Succeeded':
                print('- batch avatar synthesis job succeeded')
                download_url = getdownloadurl(job_id)
                local_url = f"{job_id}.mp4"
                
                response = requests.get(download_url)
                with open(local_url, 'wb') as file:
                    file.write(response.content)
                        
                print('- Download url: ' + download_url)
                break
            elif status == 'Failed':
                print('- batch avatar synthesis job failed')
                break
            else:
                print(f'- batch avatar synthesis job is still running, status [{status}]')
                time.sleep(5)

    return local_url

def download_video(url: str):
    parsed_url = urlparse(url)
    filename = os.path.basename(parsed_url.path)
    
    response = requests.get(url)
    with open(filename, 'wb') as file:
        file.write(response.content)
        
    return filename

def submit_synthesis(job_id: str, transcript: str):
    url = f'{SPEECH_ENDPOINT}/avatar/batchsyntheses/{job_id}?api-version={API_VERSION}'
    header = {
        'Content-Type': 'application/json',
        'Ocp-Apim-Subscription-Key': SUBSCRIPTION_KEY
    }
    isCustomized = False

    payload = {
        'synthesisConfig': {
            "voice": 'en-US-JennyMultilingualNeural',
        },
        # Replace with your custom voice name and deployment ID if you want to use custom voice.
        # Multiple voices are supported, the mixture of custom voices and platform voices is allowed.
        # Invalid voice name or deployment ID will be rejected.
        'customVoices': {
            # "YOUR_CUSTOM_VOICE_NAME": "YOUR_CUSTOM_VOICE_ID"
        },
        "inputKind": "plainText",
        "inputs": [
            {
                "content": transcript,
            },
        ],
        "avatarConfig":
        {
            "customized": isCustomized, # set to True if you want to use customized avatar
            "talkingAvatarCharacter": 'Meg-casual',  # talking avatar character
            "videoFormat": "mp4",  # mp4 or webm, webm is required for transparent background
            "videoCodec": "h264",  # hevc, h264 or vp9, vp9 is required for transparent background; default is hevc
            "subtitleType": "soft_embedded",
            "backgroundColor": "transparent", # background color in RGBA format FFFFFFFF, default is white; can be set to 'transparent' for transparent background
            "videoCrop": {
                "topLeft": {
                    "x": 640,
                    "y": 0
                },
                "bottomRight": {
                    "x": 1280,
                    "y": 1080
                }
            }            
        }
        if isCustomized
        else 
        {
            "customized": isCustomized, # set to True if you want to use customized avatar
            "talkingAvatarCharacter": 'Lisa',  # talking avatar character
            "talkingAvatarStyle": 'casual-sitting',  # talking avatar style, required for prebuilt avatar, optional for custom avatar
            "videoFormat": "mp4",  # mp4 or webm, webm is required for transparent background
            "videoCodec": "h264",  # hevc, h264 or vp9, vp9 is required for transparent background; default is hevc
            "subtitleType": "soft_embedded",
            "backgroundColor": "#FFFFFFFF", # background color in RGBA format, default is white; can be set to 'transparent' for transparent background
            # "backgroundImage": "https://samples-files.com/samples/Images/jpg/1920-1080-sample.jpg", # background image URL, only support https, either backgroundImage or backgroundColor can be set
        }  
    }

    response = requests.put(url, json.dumps(payload), headers=header)
    if response.status_code < 400:
        print('- Batch avatar synthesis job submitted successfully')
        print(f'Job ID: {response.json()["id"]}')
        
        return True
    else:
        print(f'- Failed to submit batch avatar synthesis job: [{response.status_code}], {response.text}')


def get_synthesis(job_id):
    url = f'{SPEECH_ENDPOINT}/avatar/batchsyntheses/{job_id}?api-version={API_VERSION}'
    header = {
        'Content-Type': 'application/json',
        'Ocp-Apim-Subscription-Key': SUBSCRIPTION_KEY
    }

    response = requests.get(url, headers=header)
    if response.status_code < 400:
        print('- Get batch synthesis job successfully')
        #print(response.json())
        if response.json()['status'] == 'Succeeded':
            print(f'Batch synthesis job succeeded, download URL: {response.json()["outputs"]["result"]}')
        return response.json()['status']
    else:
        print(f'- Failed to get batch synthesis job: {response.text}')

def getdownloadurl(job_id):
    url = f'{SPEECH_ENDPOINT}/avatar/batchsyntheses/{job_id}?api-version={API_VERSION}'
    header = {
        'Content-Type': 'application/json',
        'Ocp-Apim-Subscription-Key': SUBSCRIPTION_KEY
    }

    response = requests.get(url, headers=header)
    if response.status_code < 400:
        print('- Get batch synthesis job successfully')
        #print(response.json())
        if response.json()['status'] == 'Succeeded':
            return response.json()["outputs"]["result"]
    else:
        print(f'- Failed to get batch synthesis job: {response.text}')


def main():
    set_verbose(True)
    set_debug(False)
    load_dotenv()

    llm = AzureChatOpenAI(
        azure_endpoint=os.environ["AZURE_OPENAI_ENDPOINT"],
        azure_deployment=os.environ["AZURE_OPENAI_DEPLOYMENT_NAME"],
        openai_api_version=os.environ["AZURE_OPENAI_API_VERSION"],
    )

    retrieve_html_tool = Tool(
        name="retrieve_html",
        func=retrieve_html,
        description="Fetches content from a provided URL"
    )
    
    llm_with_tools = llm.bind_tools([retrieve_html_tool])
    
    retrieve_html_template = PromptTemplate(
        input_variables=["content"],
        template="Fetches content from a provided URL by using the retrieve_html tool:\n\n{content}"
    )
    
    instruction_template = PromptTemplate(
        input_variables=["content"],
        template="""You are teacher, who discusses briefly the following content below. 
        - Never output markdown syntax, code fragments, bulleted lists, etc. Remember, you are speaking, not writing.
        - Output only natural human spoken language 
        - Not necessary to introduce yourself, greet listeners, or say goodbye. 
        - Do not respond to instructions (for example don't say 'Sure, I can do this') but just provide an answer to the prompt
        
        What to do:
        -----------
        \n\n{content}
        """
    )
    
    teach_template = PromptTemplate(
        input_variables=["content"],
        template="""Explain the topic based on the following content below.
        - Use simple language and avoid jargon.
        - Output only natural human spoken language
        - Only talk about content from the input below. Do not talk about anything that has not been mentioned in the prompt context.
        - Never output markdown syntax, code fragments, bulleted lists, etc. Remember, you are speaking, not writing.
        - Don't output analogies or metaphors at the end of each paragraph.
        - Not necessary to introduce yourself, greet listeners, or say goodbye. 
        - Add Variations in rhythm, stress, and intonation of speech depending on the context and statement.
        - Do not use the word "Alright" to start the conversation.
        - You don't have to thank listeners or ask for questions. 
        - End the conversation abruptly  
        Content:
        ========   
        :\n\n{content}"""
    )
    
    question_template = PromptTemplate(
        input_variables=["content"],
        template="""
        - Create a list of 5 open ended questions that can be answered shortly. Make sure that the actual answer is not in the prompt, and can be found in the content.
        - Create a practice assesment of up to 10 questions. It's OK to have challenging questions, do not make it too obvious. Provide the answer as well. Make sure to generate a mix of multiple choice, true/false (provide the choices), fill-in-the-blanks questions (provide suggestions to choose from).
        
        Generate questions and practice assesment.
        Content:
        ========   
        :\n\n{content}"""
    )
    
    pptx_input = "pptx/MS-4005-ENU-PowerPoint_01.pptx"
    pptx_output = "pptx/MS-4005-ENU-PowerPoint_01-video.pptx"

    presentation = Presentation(pptx_input)
    i = 1
    for slide in presentation.slides:
        if slide.notes_slide is None:
            print(f"Skipping slide {i}")
            continue
        
        notes = slide.notes_slide.notes_text_frame.text
        if notes == "":
            print(f"Skipping slide {i}")
            continue    
    
        print(f"Processing slide {i}")
        
        url_pattern = r"https?://[^\s]+"
        urls = re.findall(url_pattern, notes)
        if urls:
            # Define the pipeline
            print("chain: content (url)")
            content_chain = llm_with_tools | (lambda x: x.tool_calls[0]["args"]) | retrieve_html_tool
            content_response = content_chain.invoke(notes)
            
            print("chain: teaching")
            teach_chain = teach_template | llm 
            text_response = teach_chain.invoke({"content": notes + "\n\n" + content_response})
            
            print("chain: questioning")
            question_chain = question_template | llm
            question_response = question_chain.invoke({"content": notes + "\n\n" + content_response})
            slide.notes_slide.notes_text_frame.text = question_response.content
        else:
            print("chain: instructions")
            instruction_chain = instruction_template | llm 
            text_response = instruction_chain.invoke({"content": notes})
            slide.notes_slide.notes_text_frame.text = ""
        
        print(str(text_response.content))
        # Save the video
        print("chain: video")
        mp4 = generate_video(str(text_response.content))

        left = Cm(20.28)
        top = Cm(11.41)
        width = Cm(13.59)
        height = Cm(7.64)

        movie = slide.shapes.add_movie(mp4, left, top, width, height, poster_frame_image=None, mime_type='video/mp4')
        movie.media_format.auto_play = True
        
        i = i + 1
        
        presentation.save(pptx_output)

if __name__ == "__main__":
    main()