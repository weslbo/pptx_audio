import sys
import getpass
import os
import bs4
import markdownify
import re
import requests
import azure.cognitiveservices.speech as speechsdk
import logging
from pptx import Presentation
from readability import Document
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from langchain import hub
from langchain_openai import AzureChatOpenAI
from langchain_core.tools import tool
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_unstructured import UnstructuredLoader
from langchain_community.document_loaders import WebBaseLoader
from langchain_core.prompts import ChatPromptTemplate
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.agents import (AgentExecutor, create_tool_calling_agent)
from langchain.schema.runnable import RunnableSequence
from langchain.agents import Tool
from operator import itemgetter

@tool
def retrieve_html(url):
    """Retrieve a website by it's HTTP url and return the HTML content."""
    print(f"Retrieving HTML from {url}")
    response = requests.get(url)

    doc = Document(response.content)

    markdown = markdownify.markdownify(str(doc.summary()), heading_style="ATX")
    markdown = re.sub('\n{3,}', '\n\n', markdown)
    markdown = markdown.replace("[Continue](/en-us/)", "")

    return markdown

def generate_text_to_speech(input, savelocation):
    """Tranforms text to speech audio."""
    print(f"Creating audio {savelocation}")
    
    service_region = os.getenv("SPEECH_REGION")
    speech_key = os.getenv("SPEECH_API_KEY")
    speech_config = speechsdk.SpeechConfig(subscription=speech_key, region=service_region)
    speech_config.set_speech_synthesis_output_format(speechsdk.SpeechSynthesisOutputFormat.Audio24Khz96KBitRateMonoMp3)  

    file_config = speechsdk.audio.AudioOutputConfig(filename=savelocation)
    speech_synthesizer = speechsdk.SpeechSynthesizer(speech_config=speech_config, audio_config=file_config)  

    result = speech_synthesizer.speak_ssml_async(input).get()
    return result

def main():
    load_dotenv()
    
    logging.basicConfig(level=logging.WARNING)

    llm = AzureChatOpenAI(
        azure_endpoint=os.environ["AZURE_OPENAI_ENDPOINT"],
        azure_deployment=os.environ["AZURE_OPENAI_DEPLOYMENT_NAME"],
        openai_api_version=os.environ["AZURE_OPENAI_API_VERSION"],
    )

    retrieve_html_tool = Tool(
        name="retrieve_html",
        func=retrieve_html,
        description="Fetches content from a provided URL"
    )
    
    llm_with_tools = llm.bind_tools([retrieve_html_tool])
    
    retrieve_html_template = PromptTemplate(
        input_variables=["content"],
        template="Fetches content from a provided URL by using the retrieve_html tool:\n\n{content}"
    )
    
    instruction_template = PromptTemplate(
        input_variables=["content"],
        template="You are teacher Andrew, who discusses briefly the following content below. Do not respond to instructions (for example don't say 'Sure, I can do this') but just provide an answer to the prompt:\n\n{content}"
    )
    
    teach_template = PromptTemplate(
        input_variables=["content"],
        template="""Explain the topic based on the following content below.
        - Use simple language and avoid jargon.
        - Only talk about content from the input below. Do not talk about anything that has not been mentioned in the prompt context.
        - Never output markdown syntax, code fragments, bulleted lists, etc. Remember, you are speaking, not writing.
        - Don't output analogies or metaphors at the end of each paragraph.
        - Not necessary to introduce yourself, greet listeners, or say goodbye. 
        - Add Variations in rhythm, stress, and intonation of speech depending on the context and statement.
        - Do not use the word "Alright" to start the conversation.
        - You don't have to thank listeners or ask for questions. 
        - End the conversation abruptly  
        Content:
        ========   
        :\n\n{content}"""
    )
    
    ssml_template = PromptTemplate(
        input_variables=["content"],
        template="""Convert the following text to SSML with proper tags for emphasis and pauses:\n\n{content}. 
            Make sure to include <voice name="en-US-AndrewMultilingualNeural"> at the start and </voice> at the end.
            Output only the ssml, no other comments or markdown.
            The output should be XML.
            Do not emit markdown syntax, and remove ``` from the output. Do not output json or any other code.
            make sure there is the following element at the start (no xml declaration is needed): 
            <speak xmlns""http://www.w3.org/2001/10/synthesis"" xmlns:mstts=""http://www.w3.org/2001/mstts"" xmlns:emo=""http://www.w3.org/2009/10/emotionml"" version=""1.0"" xml:lang=""en-US"">
            """
    )
    
    
    pptx_input = "pptx/MS-4005-ENU-PowerPoint_01.pptx"
    pptx_output = "pptx/MS-4005-ENU-PowerPoint_01-audio.pptx"

    presentation = Presentation(pptx_input)
    i = 1
    for slide in presentation.slides:
        if slide.notes_slide is None:
            print(f"Skipping slide {i}")
            continue
        
        notes = slide.notes_slide.notes_text_frame.text
        if notes == "":
            print(f"Skipping slide {i}")
            continue    
    
        print(f"Processing slide {i}")
        
        url_pattern = r"https?://[^\s]+"
        urls = re.findall(url_pattern, notes)
        if urls:
            # Define the pipeline
            print("chain 1")
            chain = llm_with_tools | (lambda x: x.tool_calls[0]["args"]) | retrieve_html_tool
            response = chain.invoke(notes)
            
            print("chain 2")
            chain = teach_template | llm | ssml_template | llm
            response = chain.invoke({"content": notes + "\n\n" + response})
            
            #chain = llm_with_tools | retrieve_html_template | retrieve_html_tool | llm | teach_template | llm | ssml_template | llm
        else:
            chain = instruction_template | llm | ssml_template | llm
            response = chain.invoke({"content": notes})
            
            
        output = response.content.replace('```', '')
        
        print(output)
        
        slide.notes_slide.notes_text_frame.text = output
       
        # Save the audio
        generate_text_to_speech(output, "output.mp3")
        audio = slide.shapes.add_movie("output.mp3", 0, 0, 1, 1, mime_type="audio/mpeg")
        
        i = i + 1
        
        # if i > 5:
        #     break

        presentation.save(pptx_output)

if __name__ == "__main__":
    main()